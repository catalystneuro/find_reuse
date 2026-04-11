#!/usr/bin/env python3
"""Create a scientific talk PowerPoint about DANDI data reuse analysis.

Uses CatalystNeuro brand style: navy/blue accents, Calibri font, footer bar.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path

prs = Presentation()
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT

# --- CatalystNeuro brand palette ---
CN_NAVY = RGBColor(0x02, 0x12, 0x42)
CN_BLUE = RGBColor(0x07, 0x65, 0xA5)
CN_LIGHT_BLUE = RGBColor(0x5E, 0x9B, 0xC4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF6, 0xF8)
MEDIUM_GRAY = RGBColor(0x96, 0x96, 0x96)
BODY_TEXT = RGBColor(0x32, 0x32, 0x37)

FONT = 'Calibri'
PRESENTER = 'Benjamin Dichter'
LOGO_PATH = 'assets/logo_square.png'
LOGO_HORIZ_PATH = 'assets/logo_horizontal_light.png'
BLANK_LAYOUT = prs.slide_layouts[6]

FIG_DIR = Path('output/figures')
OUT_DIR = Path('output')


# ======================================================================
# Utility functions
# ======================================================================

def set_slide_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, text, left, top, width, height,
             font_size=Pt(18), color=BODY_TEXT, bold=False,
             alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = FONT
    p.alignment = alignment
    return txBox


def add_bullets(slide, items, left, top, width, height,
                font_size=Pt(18), color=BODY_TEXT, line_spacing=Pt(20)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = font_size
        p.font.color.rgb = color
        p.font.name = FONT
        p.space_after = line_spacing
    return txBox


def add_footer(slide, section_label=''):
    bar_height = Inches(0.45)
    bar_top = SLIDE_HEIGHT - bar_height
    shape = slide.shapes.add_shape(1, Inches(0), bar_top, SLIDE_WIDTH, bar_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CN_LIGHT_BLUE
    shape.line.fill.background()
    if Path(LOGO_PATH).exists():
        from PIL import Image
        img = Image.open(LOGO_PATH)
        img_w, img_h = img.size
        logo_h = Inches(0.3)
        logo_w = int(logo_h * img_w / img_h)
        logo_top = bar_top + (bar_height - logo_h) // 2
        slide.shapes.add_picture(LOGO_PATH, Inches(0.3), int(logo_top), logo_w, int(logo_h))
    add_text(slide, f'{PRESENTER}  |  DANDI Reuse Analysis',
             Inches(1.0), bar_top, Inches(6), bar_height,
             font_size=Pt(10), color=WHITE)
    if section_label:
        txBox = add_text(slide, section_label,
                         Inches(9), bar_top, Inches(4), bar_height,
                         font_size=Pt(10), color=CN_NAVY)
        txBox.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT


def add_image_centered(slide, img_path, top=Inches(1.5), max_width=Inches(11.5),
                       max_height=Inches(5.2)):
    from PIL import Image
    img = Image.open(img_path)
    img_w, img_h = img.size
    scale_w = max_width / Emu(int(img_w * 914400 / 150))
    scale_h = max_height / Emu(int(img_h * 914400 / 150))
    scale = min(scale_w, scale_h, 1.0)
    w_emu = int(img_w * 914400 / 150 * scale)
    h_emu = int(img_h * 914400 / 150 * scale)
    left = int((SLIDE_WIDTH - w_emu) / 2)
    slide.shapes.add_picture(str(img_path), left, int(top), w_emu, h_emu)


def add_accent_line(slide, top=Inches(1.15), left=Inches(0.6), width=Inches(2.0)):
    shape = slide.shapes.add_shape(1, left, top, width, Inches(0.04))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CN_BLUE
    shape.line.fill.background()


def add_notes(slide, text):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text


# ======================================================================
# Helper: Section divider slide
# ======================================================================

def add_section_slide(title, subtitle=''):
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    set_slide_bg(slide, CN_NAVY)
    add_text(slide, title,
             Inches(0.6), Inches(2.5), Inches(12), Inches(1.5),
             font_size=Pt(40), color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.6), Inches(4.0), Inches(12), Inches(0.8),
                 font_size=Pt(20), color=CN_LIGHT_BLUE,
                 alignment=PP_ALIGN.CENTER)
    return slide


# ======================================================================
# Helper: Content slide with title + bullets
# ======================================================================

def add_content_slide(title, bullets, section='', notes=''):
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    set_slide_bg(slide)
    add_text(slide, title,
             Inches(0.6), Inches(0.4), Inches(12), Inches(0.7),
             font_size=Pt(32), color=CN_NAVY, bold=True)
    add_accent_line(slide)
    add_bullets(slide, bullets,
                left=Inches(0.8), top=Inches(1.4), width=Inches(11.5), height=Inches(5.0),
                font_size=Pt(20), line_spacing=Pt(14))
    add_footer(slide, section)
    if notes:
        add_notes(slide, notes)
    return slide


# ======================================================================
# Helper: Figure slide with title + optional subtitle
# ======================================================================

def add_figure_slide(title, img_path, subtitle='', section='', notes='',
                     max_height=Inches(5.0)):
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    set_slide_bg(slide)
    add_text(slide, title,
             Inches(0.6), Inches(0.4), Inches(12), Inches(0.7),
             font_size=Pt(32), color=CN_NAVY, bold=True)
    add_accent_line(slide)
    img_top = Inches(1.3)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.6), Inches(1.15), Inches(12), Inches(0.4),
                 font_size=Pt(16), color=MEDIUM_GRAY)
        img_top = Inches(1.6)
    add_image_centered(slide, str(img_path), top=img_top, max_height=max_height)
    add_footer(slide, section)
    if notes:
        add_notes(slide, notes)
    return slide


# ======================================================================
# SLIDE 1: Title
# ======================================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide, CN_NAVY)

add_text(slide, 'DANDI Archive\nData Reuse Analysis',
         Inches(0.6), Inches(1.8), Inches(12), Inches(2.5),
         font_size=Pt(48), color=WHITE, bold=True,
         alignment=PP_ALIGN.CENTER)

add_text(slide, 'Measuring how open neuroscience datasets generate new science',
         Inches(0.6), Inches(4.0), Inches(12), Inches(0.8),
         font_size=Pt(22), color=CN_LIGHT_BLUE,
         alignment=PP_ALIGN.CENTER)

add_text(slide, f'{PRESENTER}  \u2022  CatalystNeuro\nApril 2026',
         Inches(0.6), Inches(5.2), Inches(12), Inches(1.0),
         font_size=Pt(16), color=MEDIUM_GRAY,
         alignment=PP_ALIGN.CENTER)

if Path(LOGO_HORIZ_PATH).exists():
    from PIL import Image
    img = Image.open(LOGO_HORIZ_PATH)
    img_w, img_h = img.size
    logo_h = Inches(0.5)
    logo_w = int(logo_h * img_w / img_h)
    logo_left = int((SLIDE_WIDTH - logo_w) / 2)
    slide.shapes.add_picture(LOGO_HORIZ_PATH, logo_left, int(Inches(6.5)),
                             logo_w, int(logo_h))

add_notes(slide, "Today I will present a systematic analysis of how data deposited on the DANDI Archive gets reused by other scientists. This is one of the first comprehensive studies measuring the downstream scientific impact of an open neuroscience data repository.")


# ======================================================================
# SLIDE 2: Section — Motivation
# ======================================================================
s = add_section_slide('WHY MEASURE DATA REUSE?', 'The promise of open science requires evidence')
add_notes(s, "Before diving into methods and results, let me motivate why this question matters.")


# ======================================================================
# SLIDE 3: Motivation bullets
# ======================================================================
add_content_slide(
    'THE CASE FOR MEASURING REUSE',
    [
        '\u2013  NIH and BRAIN Initiative now mandate data sharing',
        '\u2013  DANDI Archive is the primary repository for NWB neurophysiology data',
        '\u2013  554 public, non-empty dandisets as of April 2026',
        '',
        'Key unanswered questions:',
        '    \u2022  Does deposited data actually get reused?',
        '    \u2022  How long before reuse happens?',
        '    \u2022  Who reuses data \u2014 same lab or different labs?',
        '    \u2022  How is data being reused?',
        '',
        '\u2013  Answers guide repository design, funder policy, and researcher incentives',
    ],
    section='Motivation',
    notes="Funders are investing heavily in data sharing mandates. DANDI has over 550 public datasets. But the fundamental question is whether these data are actually generating new science.",
)


# ======================================================================
# SLIDE 4: Section — Methods
# ======================================================================
s = add_section_slide('METHODS', 'An automated, LLM-assisted pipeline')
add_notes(s, "The analysis involves two main phases: linking dandisets to their primary papers, then tracking who cites those papers and classifying how they use the data.")


# ======================================================================
# SLIDE 5: Phase 1
# ======================================================================
add_figure_slide(
    'PHASE 1: DANDISET-TO-PAPER LINKAGE',
    OUT_DIR / 'dandiset_coverage_flow.png',
    subtitle='554 public dandisets \u2192 359 (65%) linked to 347 unique primary papers',
    section='Methods',
    notes="204 found via formal metadata, 155 via LLM identification with DOI validation. In total, 65% of non-empty dandisets are linked to primary papers.",
)


# ======================================================================
# SLIDE 6: Phase 2
# ======================================================================
add_figure_slide(
    'PHASE 2: CITATION ANALYSIS PIPELINE',
    OUT_DIR / 'phase2_citation_flow.png',
    subtitle='14,672 unique citing papers found, 18,827 paper-dandiset pairs classified',
    section='Methods',
    notes="For each primary paper, we query OpenAlex for all citing papers. We retrieve full text for 92% and use an LLM to classify each as REUSE, MENTION, or NEITHER. About 7% are genuine data reuse.",
)


# ======================================================================
# SLIDE 7: Paper text retrieval
# ======================================================================
add_figure_slide(
    'MULTI-SOURCE TEXT RETRIEVAL',
    OUT_DIR / 'paper_fetching_flow.png',
    subtitle='Cascading fallback chain achieves 92% full-text retrieval',
    section='Methods',
    notes="Europe PMC, PubMed Central, CrossRef, Elsevier API, Unpaywall, publisher HTML scraping, and Playwright for bioRxiv. Full text is critical because dataset references appear in methods sections, not abstracts.",
)


# ======================================================================
# SLIDE 8: How papers reference dandisets
# ======================================================================
add_figure_slide(
    'HOW PAPERS REFERENCE DANDISETS',
    OUT_DIR / 'dandiset_reference_flow.png',
    subtitle='96% cite the primary paper; only 1.5% include a direct dandiset identifier',
    section='Methods',
    notes="Most reuse papers cite the associated paper but do not formally cite the dataset itself. This means citation-based approaches like ours are essential for tracking reuse.",
)


# ======================================================================
# SLIDE 9: Section — Results
# ======================================================================
s = add_section_slide('RESULTS', 'Scale, patterns, and dynamics of data reuse')
add_notes(s, "Now let me walk you through what we found.")


# ======================================================================
# SLIDE 10: Key numbers
# ======================================================================
add_content_slide(
    'REUSE AT A GLANCE',
    [
        '554   public, non-empty dandisets analyzed',
        '359   dandisets linked to 347 unique primary papers',
        '',
        '14,672   unique citing papers screened',
        '18,827   paper-dandiset pairs classified',
        '',
        '1,306   classified as genuine data REUSE (6.9% of citations)',
        '   929   different-lab reuse (71%)',
        '   374   same-lab reuse (29%)',
        '',
        'Source: DANDI (287), Other archives (690), Unclear (326)',
    ],
    section='Results',
    notes="The headline: about 7% of papers citing dandiset primary papers represent genuine data reuse. 71% of reuse comes from different labs, which is the true measure of open science impact.",
)


# ======================================================================
# SLIDE 11: Different-lab dashboard
# ======================================================================
add_figure_slide(
    'DIFFERENT-LAB REUSE: OVERVIEW',
    FIG_DIR / 'combined_different_lab.png',
    section='Results',
    notes="6-panel overview. Panel A: source archives (Allen Institute, DANDI, CRCNS top). Panel B: top journals (bioRxiv, eLife, Nature Communications). Panel C: cumulative growth accelerating. Panel D: reuse by year, 127 in 2025. Panels E-F: MCF and instantaneous reuse rate per dandiset.",
)


# ======================================================================
# SLIDE 12: Reuse type
# ======================================================================
add_figure_slide(
    'HOW IS DATA BEING REUSED?',
    FIG_DIR / 'reuse_type.png',
    subtitle='Tool/method demonstration is the most common reuse type (25%)',
    section='Results',
    notes="Tool demos (25%), novel analysis (19%), aggregation (15%), benchmark (12%), confirmatory (12%), simulation (11%), ML training (4%). Different-lab reuse is tool-heavy; same-lab is science-heavy (novel analysis leads).",
)


# ======================================================================
# SLIDE 13: Section — Temporal Dynamics
# ======================================================================
s = add_section_slide('TEMPORAL DYNAMICS', 'When does reuse happen?')
add_notes(s, "One of the most interesting aspects is the temporal dynamics.")


# ======================================================================
# SLIDE 14: Same-lab dashboard
# ======================================================================
add_figure_slide(
    'SAME-LAB REUSE: OVERVIEW',
    FIG_DIR / 'combined_same_lab.png',
    section='Temporal Dynamics',
    notes="Same-lab reuse starts immediately and follows a saturating exponential. DANDI Archive is the top source for same-lab reuse. Novel analysis is the dominant reuse type for same-lab papers.",
)


# ======================================================================
# SLIDE 15: Section — Modeling
# ======================================================================
s = add_section_slide('MODELING & PROJECTIONS', 'Richards curve for different-lab, saturating exponential for same-lab')
add_notes(s, "We fit models to the Mean Cumulative Function to project future reuse.")


# ======================================================================
# SLIDE 16: Modeling figure
# ======================================================================
add_figure_slide(
    'REUSE DYNAMICS AND PROJECTIONS',
    FIG_DIR / 'reuse_rate_model.png',
    section='Modeling',
    notes="Panel A: MCF fits. Different-lab saturates at K=2.1 papers/dandiset (Richards), same-lab at K=2.9 (saturating exponential). Panel B: Rate peaks at 0.65/yr at year 3.9 for different-lab. Panel C: Dandiset creation growing as t^1.64. Panel D: Projected ~1,212 cumulative DANDI reuse papers by 2029.",
)


# ======================================================================
# SLIDE 17: Section — Conclusions
# ======================================================================
s = add_section_slide('CONCLUSIONS', 'What this means for open neuroscience')
add_notes(s, "Let me wrap up with key takeaways.")


# ======================================================================
# SLIDE 18: Key takeaways
# ======================================================================
add_content_slide(
    'KEY TAKEAWAYS',
    [
        '\u2013  Open neuroscience data IS being reused at scale: 1,306 reuse papers',
        '\u2013  Most reuse (71%) comes from different labs \u2014 true open science impact',
        '\u2013  Tool demonstration and novel analysis are the top reuse modes',
        '\u2013  Different-lab reuse takes ~4 years to peak \u2014 patience is needed',
        '\u2013  Each dandiset generates ~2 different-lab reuse papers over its lifetime',
        '\u2013  Projected ~1,200 cumulative DANDI reuse papers by 2029',
        '',
        '\u2013  Most reusers cite the primary paper but NOT the dataset itself',
        '    \u2192  Formal dataset citation remains rare',
    ],
    section='Conclusions',
    notes="The big picture: open data sharing is working. But funders should expect a 2-4 year lag. Repositories and journals should work harder to make dataset citation easy and expected.",
)


# ======================================================================
# SLIDE 19: Implications
# ======================================================================
add_content_slide(
    'IMPLICATIONS AND NEXT STEPS',
    [
        'For repositories:',
        '    \u2022  Improve discoverability, metadata, cross-archive linking',
        '',
        'For funders:',
        '    \u2022  Data sharing mandates are generating real downstream science',
        '',
        'For researchers:',
        '    \u2022  Depositing data leads to citations and new collaborations',
        '',
        'For journals:',
        '    \u2022  Encourage direct dataset citation (DANDI DOI)',
        '',
        'Future work:',
        '    \u2022  Extend to OpenNeuro, CRCNS, and other archives',
        '    \u2022  Investigate what makes datasets more reusable',
        '    \u2022  Build a real-time reuse tracking dashboard',
    ],
    section='Conclusions',
    notes="The implications span all stakeholders. Individual researchers should be encouraged that sharing data leads to real impact.",
)


# ======================================================================
# SLIDE 20: Closing
# ======================================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide, CN_NAVY)

add_text(slide, 'THANK YOU',
         Inches(0.6), Inches(2.0), Inches(12), Inches(1.5),
         font_size=Pt(48), color=WHITE, bold=True,
         alignment=PP_ALIGN.CENTER)

add_text(slide, 'Benjamin Dichter  \u2022  CatalystNeuro\nben.dichter@catalystneuro.com',
         Inches(0.6), Inches(3.5), Inches(12), Inches(1.0),
         font_size=Pt(20), color=CN_LIGHT_BLUE,
         alignment=PP_ALIGN.CENTER)

add_text(slide, 'DANDI Archive: dandiarchive.org\nCode & data: github.com/catalystneuro/find_reuse',
         Inches(0.6), Inches(4.8), Inches(12), Inches(1.0),
         font_size=Pt(16), color=MEDIUM_GRAY,
         alignment=PP_ALIGN.CENTER)

if Path(LOGO_HORIZ_PATH).exists():
    from PIL import Image
    img = Image.open(LOGO_HORIZ_PATH)
    img_w, img_h = img.size
    logo_h = Inches(0.5)
    logo_w = int(logo_h * img_w / img_h)
    logo_left = int((SLIDE_WIDTH - logo_w) / 2)
    slide.shapes.add_picture(LOGO_HORIZ_PATH, logo_left, int(Inches(6.2)),
                             logo_w, int(logo_h))

add_notes(slide, "Thank you for your attention. The code, data, and analysis are all openly available.")


# ======================================================================
# Save
# ======================================================================
output_path = 'output/dandi_reuse_talk.pptx'
prs.save(output_path)
print(f'Saved {len(prs.slides)} slides to {output_path}')
