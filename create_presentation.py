"""Create a PowerPoint presentation combining all DANDI reuse analysis diagrams and plots.

Uses CatalystNeuro brand style: white background, Calibri font, navy/blue accents,
footer bar with logo.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path

prs = Presentation()
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT

# --- CatalystNeuro brand palette ---
CN_NAVY = RGBColor(0x02, 0x12, 0x42)
CN_BLUE = RGBColor(0x07, 0x65, 0xA5)
CN_LIGHT_BLUE = RGBColor(0x5E, 0x9B, 0xC4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF6, 0xF8)
MEDIUM_GRAY = RGBColor(0x96, 0x96, 0x96)
BODY_TEXT = RGBColor(0x32, 0x32, 0x37)
BLACK = RGBColor(0x00, 0x00, 0x00)

FONT = 'Calibri'
PRESENTER = 'Benjamin Dichter'
PRES_TITLE = 'DANDI Reuse Analysis'
LOGO_PATH = 'assets/logo_square.png'
LOGO_HORIZ_PATH = 'assets/logo_horizontal_light.png'

BLANK_LAYOUT = prs.slide_layouts[6]  # blank


# ======================================================================
# Utility functions
# ======================================================================

def set_slide_bg(slide, color=WHITE):
    """Set slide background to a solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, text, left, top, width, height,
             font_size=Pt(18), color=BODY_TEXT, bold=False,
             alignment=PP_ALIGN.LEFT, font=FONT):
    """Add a text box with specified properties."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = alignment
    return txBox


def add_bullets(slide, items, left, top, width, height,
                font_size=Pt(16), color=BODY_TEXT, line_spacing=Pt(24)):
    """Add a bulleted list."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = font_size
        p.font.color.rgb = color
        p.font.name = FONT
        p.space_after = line_spacing
    return txBox


def add_footer(slide, section_label=''):
    """Add CatalystNeuro footer bar with logo, presenter, and section label."""
    bar_height = Inches(0.45)
    bar_top = SLIDE_HEIGHT - bar_height

    # Footer bar background
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(0), bar_top, SLIDE_WIDTH, bar_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = CN_NAVY
    shape.line.fill.background()

    # Logo in footer (left side)
    logo_h = Inches(0.3)
    logo_top = bar_top + (bar_height - logo_h) // 2
    if Path(LOGO_PATH).exists():
        from PIL import Image
        img = Image.open(LOGO_PATH)
        img_w, img_h = img.size
        logo_w = int(logo_h * img_w / img_h)
        slide.shapes.add_picture(LOGO_PATH, Inches(0.3), int(logo_top),
                                 logo_w, int(logo_h))

    # Presenter name (center-left)
    add_text(slide, f'{PRESENTER}  |  {PRES_TITLE}',
             Inches(1.0), bar_top, Inches(6), bar_height,
             font_size=Pt(10), color=WHITE, bold=False)

    # Section label (right side)
    if section_label:
        txBox = add_text(slide, section_label,
                         Inches(9), bar_top, Inches(4), bar_height,
                         font_size=Pt(10), color=CN_LIGHT_BLUE, bold=False)
        txBox.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT


def add_image_centered(slide, img_path, top=Inches(1.5), max_width=Inches(11.5),
                       max_height=Inches(5.2)):
    """Add an image centered horizontally, fitting within max bounds."""
    from PIL import Image
    img = Image.open(img_path)
    img_w, img_h = img.size

    # Calculate scale to fit within bounds (images are 150 dpi)
    scale_w = max_width / Emu(int(img_w * 914400 / 150))
    scale_h = max_height / Emu(int(img_h * 914400 / 150))
    scale = min(scale_w, scale_h, 1.0)  # don't upscale

    w_emu = int(img_w * 914400 / 150 * scale)
    h_emu = int(img_h * 914400 / 150 * scale)

    # Center horizontally
    left = int((SLIDE_WIDTH - w_emu) / 2)

    slide.shapes.add_picture(str(img_path), left, int(top), w_emu, h_emu)


def add_accent_line(slide, top=Inches(1.15), left=Inches(0.6), width=Inches(2.0)):
    """Add a thin accent line under the title."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        left, top, width, Inches(0.04)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = CN_BLUE
    shape.line.fill.background()


# ======================================================================
# Slide 1: Title slide
# ======================================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide, CN_NAVY)

# Large centered title
add_text(slide, 'DANDI Archive\nData Reuse Analysis',
         Inches(0.6), Inches(1.8), Inches(12), Inches(2.5),
         font_size=Pt(48), color=WHITE, bold=True,
         alignment=PP_ALIGN.CENTER)

# Subtitle
add_text(slide, 'Tracking how neuroscience datasets are cited and reused',
         Inches(0.6), Inches(4.0), Inches(12), Inches(0.8),
         font_size=Pt(22), color=CN_LIGHT_BLUE, bold=False,
         alignment=PP_ALIGN.CENTER)

# Author & date
add_text(slide, f'{PRESENTER}  \u2022  CatalystNeuro\nMarch 2026',
         Inches(0.6), Inches(5.2), Inches(12), Inches(1.0),
         font_size=Pt(16), color=MEDIUM_GRAY, bold=False,
         alignment=PP_ALIGN.CENTER)

# Logo on title slide (bottom-center)
if Path(LOGO_HORIZ_PATH).exists():
    from PIL import Image
    img = Image.open(LOGO_HORIZ_PATH)
    img_w, img_h = img.size
    logo_h = Inches(0.5)
    logo_w = int(logo_h * img_w / img_h)
    logo_left = int((SLIDE_WIDTH - logo_w) / 2)
    slide.shapes.add_picture(LOGO_HORIZ_PATH, logo_left, int(Inches(6.5)),
                             logo_w, int(logo_h))


# ======================================================================
# Slide 2: Key Numbers
# ======================================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide)

add_text(slide, 'Key Numbers',
         Inches(0.6), Inches(0.4), Inches(12), Inches(0.7),
         font_size=Pt(32), color=CN_NAVY, bold=True)
add_accent_line(slide)

add_bullets(slide, [
    '792  dandisets on DANDI Archive',
    '216  dandisets with \u2265 1 primary paper (27%)',
    '262  unique primary paper DOIs (including preprint \u2194 published alternates)',
    '9,910  papers cite those primary papers after dandiset creation',
    '',
    'Paper source breakdown:',
    '  \u2022  141 dandisets: primary paper in relatedResource metadata',
    '  \u2022   80 dandisets: primary paper DOI found in description',
    '  \u2022     5 dandisets: found via both sources',
], left=Inches(0.8), top=Inches(1.5), width=Inches(11.5), height=Inches(5.0),
   font_size=Pt(20), line_spacing=Pt(12))

add_footer(slide, 'Overview')


# ======================================================================
# Slide 3: Citation Discovery Pipeline
# ======================================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide)

add_text(slide, 'Citation Discovery Pipeline',
         Inches(0.6), Inches(0.4), Inches(12), Inches(0.7),
         font_size=Pt(32), color=CN_NAVY, bold=True)
add_accent_line(slide)

add_image_centered(slide, 'output/citation_pipeline_flow.png',
                   top=Inches(1.2), max_height=Inches(5.5))

add_footer(slide, 'Pipeline')


# ======================================================================
# Slide 4: Search Reference Flow (NEW)
# ======================================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide)

add_text(slide, 'Discovering Papers That Reference Datasets',
         Inches(0.6), Inches(0.4), Inches(12), Inches(0.7),
         font_size=Pt(32), color=CN_NAVY, bold=True)
add_accent_line(slide)

add_text(slide, 'Search engines + text pattern matching across multiple archives',
         Inches(0.6), Inches(1.1), Inches(12), Inches(0.5),
         font_size=Pt(16), color=MEDIUM_GRAY)

add_image_centered(slide, 'output/search_reference_flow.png',
                   top=Inches(1.5), max_height=Inches(5.2))

add_footer(slide, 'Pipeline')


# ======================================================================
# Slide 5: Paper Fetching Flow
# ======================================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide)

add_text(slide, 'Paper Text Fetching Flow',
         Inches(0.6), Inches(0.4), Inches(12), Inches(0.7),
         font_size=Pt(32), color=CN_NAVY, bold=True)
add_accent_line(slide)

add_text(slide, 'Multi-source fallback chain for retrieving full paper text',
         Inches(0.6), Inches(1.1), Inches(12), Inches(0.5),
         font_size=Pt(16), color=MEDIUM_GRAY)

add_image_centered(slide, 'output/paper_fetching_flow.png',
                   top=Inches(1.5), max_height=Inches(5.2))

add_footer(slide, 'Pipeline')


# ======================================================================
# Slide 6: How Papers Reference Dandisets
# ======================================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide)

add_text(slide, 'How Papers Reference Dandisets',
         Inches(0.6), Inches(0.4), Inches(12), Inches(0.7),
         font_size=Pt(32), color=CN_NAVY, bold=True)
add_accent_line(slide)

add_image_centered(slide, 'output/dandiset_reference_flow.png',
                   top=Inches(1.3), max_height=Inches(5.3))

add_footer(slide, 'Results')


# ======================================================================
# Slide 7: Cumulative Citations Over Time
# ======================================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide)

add_text(slide, 'Cumulative Dataset Citations Over Time',
         Inches(0.6), Inches(0.4), Inches(12), Inches(0.7),
         font_size=Pt(32), color=CN_NAVY, bold=True)
add_accent_line(slide)

add_image_centered(slide, 'dandi_citations_quarterly.png',
                   top=Inches(1.3), max_height=Inches(5.3))

add_footer(slide, 'Results')


# ======================================================================
# Slide 8: Time to Reuse - from dandiset creation
# ======================================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide)

add_text(slide, 'Time from Dandiset Creation to Secondary Publication',
         Inches(0.6), Inches(0.4), Inches(12), Inches(0.7),
         font_size=Pt(32), color=CN_NAVY, bold=True)
add_accent_line(slide)

add_image_centered(slide, 'output/time_to_reuse_histogram.png',
                   top=Inches(1.3), max_height=Inches(5.3))

add_footer(slide, 'Results')


# ======================================================================
# Slide 9: Time to Reuse - from primary paper
# ======================================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide)

add_text(slide, 'Time from Primary Paper to Secondary Publication',
         Inches(0.6), Inches(0.4), Inches(12), Inches(0.7),
         font_size=Pt(32), color=CN_NAVY, bold=True)
add_accent_line(slide)

add_image_centered(slide, 'output/time_to_reuse_histogram_from_primary.png',
                   top=Inches(1.3), max_height=Inches(5.3))

add_footer(slide, 'Results')


# NOTE: Kaplan-Meier survival analysis slide removed per user request
# (it only counts first use of each dandiset)


# ======================================================================
# Slide 10: MCF - Reuse Papers per Dandiset
# ======================================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide)

add_text(slide, 'Mean Cumulative Function: Reuse Papers per Dandiset',
         Inches(0.6), Inches(0.4), Inches(12), Inches(0.7),
         font_size=Pt(28), color=CN_NAVY, bold=True)
add_accent_line(slide)

add_text(slide, 'Same lab vs. different lab reuse over time',
         Inches(0.6), Inches(1.1), Inches(12), Inches(0.5),
         font_size=Pt(16), color=MEDIUM_GRAY)

add_image_centered(slide, 'output/reuse_prediction_mcf.png',
                   top=Inches(1.5), max_height=Inches(5.2))

add_footer(slide, 'Predictions')


# ======================================================================
# Slide 11: Predicted Cumulative Dandiset Reuse
# ======================================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide)

add_text(slide, 'Observed and Predicted Cumulative Dandiset Reuse',
         Inches(0.6), Inches(0.4), Inches(12), Inches(0.7),
         font_size=Pt(28), color=CN_NAVY, bold=True)
add_accent_line(slide)

add_image_centered(slide, 'output/reuse_survival_analysis_prediction.png',
                   top=Inches(1.3), max_height=Inches(5.3))

add_footer(slide, 'Predictions')


# ======================================================================
# Slide 12: Predicted Cumulative Reuse Papers
# ======================================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide)

add_text(slide, 'Observed and Predicted Cumulative Reuse Papers',
         Inches(0.6), Inches(0.4), Inches(12), Inches(0.7),
         font_size=Pt(28), color=CN_NAVY, bold=True)
add_accent_line(slide)

add_text(slide, 'Same lab vs. different lab, with and without new dandiset growth',
         Inches(0.6), Inches(1.1), Inches(12), Inches(0.5),
         font_size=Pt(16), color=MEDIUM_GRAY)

add_image_centered(slide, 'output/reuse_prediction.png',
                   top=Inches(1.5), max_height=Inches(5.2))

add_footer(slide, 'Predictions')


# ======================================================================
# Save
# ======================================================================
output_path = 'output/dandi_reuse_analysis.pptx'
prs.save(output_path)
print(f"Saved presentation to {output_path}")
print(f"  {len(prs.slides)} slides")
